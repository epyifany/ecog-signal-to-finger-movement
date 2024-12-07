from pptx import Presentation

# Sample text data
content = {
    "Slide 1 Title": "This is the first slide",
    "Slide 2 Title": "Here's another slide",
    "Slide 3 Title": "Wrapping up"
}

# Create a PowerPoint presentation
presentation = Presentation()

for title, body in content.items():
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body

# Save the presentation
presentation.save("output.pptx")
